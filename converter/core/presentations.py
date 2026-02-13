"""PPTX <-> PDF conversion."""

from pathlib import Path

from converter.utils.helpers import get_output_path, validate_input_file


def pptx_to_pdf(input_path: str, output_dir: str | None = None) -> Path:
    """Convert a PPTX file to PDF using Microsoft PowerPoint (COM automation)."""
    import win32com.client
    import pythoncom

    input_file = validate_input_file(input_path)
    output_path = get_output_path(input_path, "pdf", output_dir)

    pythoncom.CoInitialize()
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            str(input_file.resolve()), WithWindow=False
        )
        # 32 = ppSaveAsPDF
        presentation.SaveAs(str(output_path.resolve()), 32)
        presentation.Close()
        powerpoint.Quit()
    finally:
        pythoncom.CoUninitialize()

    return output_path


def pdf_to_pptx(input_path: str, output_dir: str | None = None) -> Path:
    """Convert a PDF file to PPTX (each page becomes a slide with the page as an image)."""
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    import tempfile
    import os

    input_file = validate_input_file(input_path)
    output_path = get_output_path(input_path, "pptx", output_dir)

    # Convert PDF pages to images
    images = convert_from_path(str(input_file))

    # Create a presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    with tempfile.TemporaryDirectory() as tmp_dir:
        for i, image in enumerate(images):
            # Save page image temporarily
            img_path = os.path.join(tmp_dir, f"page_{i}.png")
            image.save(img_path, "PNG")

            # Add slide with the image
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_path,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

    prs.save(str(output_path))
    return output_path
